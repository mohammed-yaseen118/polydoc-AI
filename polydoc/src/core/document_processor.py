"""
PolyDoc - Multi-format Document Processor
Handles PDF, DOCX, PPTX, and image files with layout preservation
"""

import os
import asyncio
from typing import Dict, List, Any, Optional, Tuple
from dataclasses import dataclass
from pathlib import Path
import logging

# Document processing imports
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
import cv2
import numpy as np
import pandas as pd
import json
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
import markdown
import chardet

# OCR imports
import pytesseract
import easyocr

# Additional imports for new formats
try:
    import openpyxl
    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False

try:
    from odf.opendocument import load
    from odf.text import P
    from odf import teletype
    ODT_AVAILABLE = True
except ImportError:
    ODT_AVAILABLE = False

# Layout analysis (optional)
try:
    import layoutparser as lp
    LAYOUT_PARSER_AVAILABLE = True
except ImportError:
    lp = None
    LAYOUT_PARSER_AVAILABLE = False

@dataclass
class DocumentElement:
    """Represents a document element with layout information"""
    text: str
    page_number: int
    element_type: str  # 'paragraph', 'heading', 'table', 'image', 'handwriting'
    bbox: Tuple[float, float, float, float]  # (x1, y1, x2, y2)
    confidence: float
    language: Optional[str] = None
    font_info: Optional[Dict] = None

@dataclass
class ProcessedDocument:
    """Container for processed document information"""
    filename: str
    total_pages: int
    elements: List[DocumentElement]
    summary: Optional[str] = None
    metadata: Dict[str, Any] = None

class DocumentProcessor:
    """Main document processor supporting multiple formats"""
    
    def __init__(self):
        self.logger = logging.getLogger(__name__)
        
        # Initialize EasyOCR with smart language selection and memory optimization
        # EasyOCR has compatibility restrictions between certain language combinations
        self.ocr_readers = {}
        self.primary_ocr_reader = None
        self.use_tesseract_fallback = False
        
        # Define compatible language groups for EasyOCR (starting with most important)
        # Note: Some languages may not be compatible together, so we test individual pairs
        compatible_groups = [
            ['en'],         # Start with English only for stability
            ['en', 'hi'],   # English + Hindi (Devanagari)
            ['en', 'kn'],   # English + Kannada
            ['hi'],         # Hindi only
            ['kn'],         # Kannada only
        ]
        
        # Try to initialize OCR readers for each compatible group with better error handling
        initialized_count = 0
        for i, languages in enumerate(compatible_groups):
            try:
                self.logger.info(f"Attempting to initialize OCR with languages: {languages}")
                # Use more conservative settings to avoid memory issues
                reader = easyocr.Reader(
                    languages, 
                    gpu=False,  # Force CPU to avoid GPU issues
                    verbose=False,  # Reduce logging
                    download_enabled=True  # Allow downloading if needed
                )
                group_key = '+'.join(languages)
                self.ocr_readers[group_key] = reader
                
                # Set the first successful reader as primary
                if self.primary_ocr_reader is None:
                    self.primary_ocr_reader = reader
                    self.logger.info(f"✅ Primary OCR initialized with: {languages}")
                else:
                    self.logger.info(f"✅ Additional OCR reader initialized: {languages}")
                
                initialized_count += 1
                
                # For efficiency and memory, limit to fewer combinations
                if initialized_count >= 2:  # Reduced from 3 to 2
                    self.logger.info(f"Initialized {initialized_count} OCR readers, stopping for efficiency")
                    break
                
            except Exception as e:
                error_msg = str(e)
                if "not enough memory" in error_msg.lower() or "allocate" in error_msg.lower():
                    self.logger.error(f"❌ Memory allocation failed for {languages}: {e}")
                    self.logger.info("Switching to Tesseract fallback due to memory constraints...")
                    self.use_tesseract_fallback = True
                    break
                else:
                    self.logger.warning(f"❌ Could not initialize OCR for {languages}: {e}")
                continue
        
        if initialized_count == 0 or self.use_tesseract_fallback:
            # Fallback to Tesseract when EasyOCR fails
            try:
                self.logger.info("Attempting Tesseract OCR fallback initialization...")
                # Test if tesseract is available
                import pytesseract
                pytesseract.image_to_string(Image.new('RGB', (100, 30), color='white'))
                self.use_tesseract_fallback = True
                self.primary_ocr_reader = None  # Will use tesseract
                self.logger.info("✅ Tesseract OCR fallback initialized")
                initialized_count = 1
            except Exception as e:
                self.logger.error(f"❌ Failed to initialize Tesseract OCR: {e}")
                self.logger.warning("⚠️ No OCR available - image processing will be limited")
                self.primary_ocr_reader = None
                self.use_tesseract_fallback = False
        else:
            self.logger.info(f"✅ Successfully initialized {initialized_count} OCR reader groups")
        
        # For backward compatibility, set the main ocr_reader attribute
        self.ocr_reader = self.primary_ocr_reader
        
        # Initialize layout model if layoutparser is available
        self.layout_model = None
        if LAYOUT_PARSER_AVAILABLE:
            try:
                self.layout_model = lp.Detectron2LayoutModel(
                    'lp://PubLayNet/faster_rcnn_R_50_FPN_3x/config',
                    extra_config=["MODEL.ROI_HEADS.SCORE_THRESH_TEST", 0.8],
                    label_map={0: "Text", 1: "Title", 2: "List", 3: "Table", 4: "Figure"}
                )
            except Exception as e:
                self.logger.warning(f"Layout model not available: {e}")
                self.layout_model = None
        else:
            self.logger.info("Layout parser not available, skipping layout model initialization")
        
        # Supported file types (expanded)
        self.supported_formats = {
            # Original formats
            '.pdf', '.docx', '.ppt', '.pptx', '.png', '.jpg', '.jpeg', '.tiff', '.bmp',
            # New text formats
            '.txt', '.rtf', '.md', '.markdown',
            # Spreadsheet formats
            '.csv', '.xlsx', '.xls',
            # Structured data formats
            '.json', '.xml', '.html', '.htm',
            # OpenDocument formats
            '.odt'
        }
    
    def _get_best_ocr_reader(self, detected_language: str = None):
        """Get the best OCR reader for the detected language"""
        if not detected_language or detected_language == 'en':
            return self.primary_ocr_reader
        
        # Try to find an OCR reader that includes the detected language
        for reader_key, reader in self.ocr_readers.items():
            if detected_language in reader_key:
                return reader
        
        # If no specific reader found, return primary
        return self.primary_ocr_reader
    
    def estimate_processing_time(self, file_path: str) -> dict:
        """Estimate processing time based on file size and type"""
        try:
            file_path = Path(file_path)
            file_size_mb = file_path.stat().st_size / (1024 * 1024)  # Size in MB
            file_ext = file_path.suffix.lower()
            
            # Base processing times (seconds per MB)
            time_estimates = {
                '.pdf': 3,      # PDF processing is moderate
                '.docx': 1,     # DOCX is fastest
                '.ppt': 4,      # PPT processing requires conversion
                '.pptx': 2,     # PPTX is moderate
                '.png': 8,      # Images need OCR - slower
                '.jpg': 8,
                '.jpeg': 8,
                '.tiff': 10,    # TIFF can be large
                '.bmp': 12      # BMP is uncompressed - slowest
            }
            
            base_time = time_estimates.get(file_ext, 8)  # Default to 8s/MB
            estimated_seconds = max(5, int(file_size_mb * base_time))  # Minimum 5 seconds
            
            return {
                'estimated_seconds': estimated_seconds,
                'estimated_minutes': round(estimated_seconds / 60, 1),
                'file_size_mb': round(file_size_mb, 2),
                'complexity': 'High' if estimated_seconds > 60 else 'Medium' if estimated_seconds > 20 else 'Low'
            }
        except Exception as e:
            self.logger.error(f"Error estimating time: {e}")
            return {
                'estimated_seconds': 30,
                'estimated_minutes': 0.5,
                'file_size_mb': 0,
                'complexity': 'Unknown'
            }
    
    async def process_document(self, file_path: str) -> ProcessedDocument:
        """Main method to process any supported document format"""
        try:
            file_path = Path(file_path)
            
            if not file_path.exists():
                raise FileNotFoundError(f"File not found: {file_path}")
            
            if file_path.suffix.lower() not in self.supported_formats:
                raise ValueError(f"Unsupported format: {file_path.suffix}")
            
            self.logger.info(f"Processing document: {file_path.name}")
            
            # Route to appropriate processor
            file_ext = file_path.suffix.lower()
            
            if file_ext == '.pdf':
                return await self._process_pdf(file_path)
            elif file_ext == '.docx':
                return await self._process_docx(file_path)
            elif file_ext == '.ppt':
                return await self._process_ppt(file_path)
            elif file_ext == '.pptx':
                return await self._process_pptx(file_path)
            elif file_ext in {'.txt', '.rtf'}:
                return await self._process_text(file_path)
            elif file_ext in {'.md', '.markdown'}:
                return await self._process_markdown(file_path)
            elif file_ext in {'.csv', '.xlsx', '.xls'}:
                return await self._process_spreadsheet(file_path)
            elif file_ext == '.json':
                return await self._process_json(file_path)
            elif file_ext == '.xml':
                return await self._process_xml(file_path)
            elif file_ext in {'.html', '.htm'}:
                return await self._process_html(file_path)
            elif file_ext == '.odt':
                return await self._process_odt(file_path)
            else:  # Image formats
                return await self._process_image(file_path)
        
        except Exception as e:
            self.logger.error(f"Error processing document: {e}")
            raise
    
    async def _process_pdf(self, file_path: Path) -> ProcessedDocument:
        """Process PDF files with layout preservation"""
        elements = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            total_pages = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                # Extract text
                text_content = page.extract_text()
                
                if text_content.strip():
                    # For PDFs, we'll treat the extracted text as paragraphs
                    paragraphs = [p.strip() for p in text_content.split('\n\n') if p.strip()]
                    
                    for i, paragraph in enumerate(paragraphs):
                        element = DocumentElement(
                            text=paragraph,
                            page_number=page_num,
                            element_type='paragraph',
                            bbox=(0, i*50, 500, (i+1)*50),  # Approximate positioning
                            confidence=0.9,
                            language=self._detect_language(paragraph)
                        )
                        elements.append(element)
                
                # Convert PDF page to image for layout analysis if needed
                # This would require pdf2image library for more advanced layout detection
        
        return ProcessedDocument(
            filename=file_path.name,
            total_pages=total_pages,
            elements=elements,
            metadata={'file_type': 'pdf', 'size': file_path.stat().st_size}
        )
    
    async def _process_docx(self, file_path: Path) -> ProcessedDocument:
        """Process DOCX files with enhanced content extraction"""
        try:
            doc = DocxDocument(file_path)
            elements = []
            page_number = 1
            
            self.logger.info(f"Processing DOCX with {len(doc.paragraphs)} paragraphs and {len(doc.tables)} tables")
            
            # Process paragraphs with better error handling
            for para_idx, paragraph in enumerate(doc.paragraphs):
                if paragraph.text and paragraph.text.strip():
                    try:
                        # Determine element type based on style
                        style_name = paragraph.style.name if paragraph.style else 'Normal'
                        element_type = 'heading' if 'heading' in style_name.lower() else 'paragraph'
                        
                        # Clean and normalize text
                        cleaned_text = paragraph.text.strip().replace('\r\n', '\n').replace('\r', '\n')
                        
                        if len(cleaned_text) > 5:  # Only add meaningful content
                            element = DocumentElement(
                                text=cleaned_text,
                                page_number=page_number,
                                element_type=element_type,
                                bbox=(0, para_idx*30, 500, (para_idx+1)*30),
                                confidence=1.0,
                                language=self._detect_language(cleaned_text),
                                font_info={'style': style_name}
                            )
                            elements.append(element)
                    except Exception as para_error:
                        self.logger.warning(f"Error processing paragraph {para_idx}: {para_error}")
                        continue
            
            # Process tables with better error handling
            for table_idx, table in enumerate(doc.tables):
                try:
                    table_rows = []
                    for row in table.rows:
                        row_cells = []
                        for cell in row.cells:
                            cell_text = cell.text.strip() if cell.text else ""
                            row_cells.append(cell_text)
                        if any(cell.strip() for cell in row_cells):  # Only add non-empty rows
                            table_rows.append(' | '.join(row_cells))
                    
                    if table_rows:
                        table_content = '\n'.join(table_rows)
                        element = DocumentElement(
                            text=table_content,
                            page_number=page_number,
                            element_type='table',
                            bbox=(0, 1000+table_idx*100, 500, 1000+(table_idx+1)*100),
                            confidence=1.0,
                            language=self._detect_language(table_content)
                        )
                        elements.append(element)
                except Exception as table_error:
                    self.logger.warning(f"Error processing table {table_idx}: {table_error}")
                    continue
            
            # Extract text from headers and footers
            try:
                for section in doc.sections:
                    # Process header
                    if section.header:
                        for paragraph in section.header.paragraphs:
                            if paragraph.text and paragraph.text.strip():
                                element = DocumentElement(
                                    text=paragraph.text.strip(),
                                    page_number=page_number,
                                    element_type='header',
                                    bbox=(0, -50, 500, 0),
                                    confidence=0.9,
                                    language=self._detect_language(paragraph.text)
                                )
                                elements.append(element)
                    
                    # Process footer
                    if section.footer:
                        for paragraph in section.footer.paragraphs:
                            if paragraph.text and paragraph.text.strip():
                                element = DocumentElement(
                                    text=paragraph.text.strip(),
                                    page_number=page_number,
                                    element_type='footer',
                                    bbox=(0, 1100, 500, 1150),
                                    confidence=0.9,
                                    language=self._detect_language(paragraph.text)
                                )
                                elements.append(element)
            except Exception as header_footer_error:
                self.logger.warning(f"Error processing headers/footers: {header_footer_error}")
            
            # If no elements were extracted, create a fallback element
            if not elements:
                self.logger.warning(f"No content extracted from DOCX file: {file_path.name}")
                fallback_element = DocumentElement(
                    text=f"Document processed but no readable content found in {file_path.name}. The file may be empty, corrupted, or contain only images/formatting.",
                    page_number=1,
                    element_type='error',
                    bbox=(0, 0, 500, 50),
                    confidence=0.5,
                    language='en'
                )
                elements.append(fallback_element)
            
            self.logger.info(f"Successfully extracted {len(elements)} elements from DOCX")
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=elements,
                metadata={
                    'file_type': 'docx', 
                    'size': file_path.stat().st_size,
                    'paragraphs_count': len([e for e in elements if e.element_type == 'paragraph']),
                    'tables_count': len([e for e in elements if e.element_type == 'table']),
                    'headings_count': len([e for e in elements if e.element_type == 'heading'])
                }
            )
            
        except Exception as e:
            self.logger.error(f"Error processing DOCX file {file_path}: {e}")
            # Create error document
            error_element = DocumentElement(
                text=f"Error processing DOCX file: {str(e)}. Please ensure the file is not corrupted and try again.",
                page_number=1,
                element_type='error',
                bbox=(0, 0, 500, 50),
                confidence=0.0,
                language='en'
            )
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=[error_element],
                metadata={'file_type': 'docx', 'error': str(e)}
            )
    
    async def _process_ppt(self, file_path: Path) -> ProcessedDocument:
        """Process legacy PPT files - fallback approach"""
        try:
            self.logger.info(f"Processing legacy PPT file: {file_path.name}")
            
            # Legacy PPT files are not directly supported by python-pptx
            # We'll create a placeholder document with a message about conversion needed
            elements = []
            
            # Add informational element about PPT processing
            conversion_info = DocumentElement(
                text=f"Legacy PowerPoint file (.ppt) detected: {file_path.name}. "
                      "For best results, please convert to .pptx format. "
                      "Currently showing basic file information only.",
                page_number=1,
                element_type='heading',
                bbox=(0, 0, 500, 100),
                confidence=1.0,
                language='en',
                font_info={'type': 'conversion_notice'}
            )
            elements.append(conversion_info)
            
            # Try to extract basic file information
            try:
                file_size = file_path.stat().st_size
                size_mb = file_size / (1024 * 1024)
                
                file_info = DocumentElement(
                    text=f"File Information:\n"
                          f"- Filename: {file_path.name}\n"
                          f"- Size: {size_mb:.2f} MB\n"
                          f"- Format: Legacy PowerPoint (.ppt)\n"
                          f"- Recommendation: Convert to .pptx for full text extraction",
                    page_number=1,
                    element_type='paragraph',
                    bbox=(0, 100, 500, 300),
                    confidence=1.0,
                    language='en',
                    font_info={'type': 'file_info'}
                )
                elements.append(file_info)
                
            except Exception as info_error:
                self.logger.warning(f"Could not extract file info from PPT: {info_error}")
            
            # Alternative approach: Try to read as binary and look for text strings
            try:
                text_content = self._extract_text_from_binary_ppt(file_path)
                if text_content:
                    for i, text in enumerate(text_content):
                        if text.strip():
                            element = DocumentElement(
                                text=text.strip(),
                                page_number=1,
                                element_type='paragraph',
                                bbox=(0, 300 + i*30, 500, 330 + i*30),
                                confidence=0.5,  # Lower confidence for binary extraction
                                language=self._detect_language(text),
                                font_info={'extraction_method': 'binary_strings'}
                            )
                            elements.append(element)
            except Exception as binary_error:
                self.logger.warning(f"Binary text extraction failed: {binary_error}")
            
            metadata = {
                'file_type': 'ppt',
                'size': file_path.stat().st_size,
                'processing_note': 'Legacy PPT format - limited extraction capabilities',
                'recommendation': 'Convert to PPTX for full feature support'
            }
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,  # Cannot determine actual slide count from binary
                elements=elements,
                metadata=metadata
            )
            
        except Exception as e:
            self.logger.error(f"Error processing PPT file: {e}")
            # Return a basic error document
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=[DocumentElement(
                    text=f"Failed to process PPT file {file_path.name}. Error: {str(e)}. "
                          "Please try converting to PPTX format.",
                    page_number=1,
                    element_type='error',
                    bbox=(0, 0, 500, 100),
                    confidence=0.0,
                    language='en'
                )],
                metadata={'file_type': 'ppt', 'error': str(e), 'size': file_path.stat().st_size if file_path.exists() else 0}
            )
    
    async def _process_pptx(self, file_path: Path) -> ProcessedDocument:
        """Enhanced PowerPoint processing with comprehensive content extraction"""
        try:
            prs = Presentation(file_path)
            elements = []
            
            self.logger.info(f"Processing PowerPoint with {len(prs.slides)} slides")
            
            for slide_num, slide in enumerate(prs.slides, 1):
                # Extract slide title first (if exists)
                slide_title = self._extract_slide_title(slide)
                if slide_title:
                    elements.append(DocumentElement(
                        text=slide_title,
                        page_number=slide_num,
                        element_type='heading',
                        bbox=(0, 0, 500, 50),
                        confidence=1.0,
                        language=self._detect_language(slide_title),
                        font_info={'type': 'slide_title', 'slide_number': slide_num}
                    ))
                
                # Process all shapes in the slide
                for shape_idx, shape in enumerate(slide.shapes):
                    try:
                        # Handle text shapes
                        if hasattr(shape, 'text') and shape.text.strip():
                            # Enhanced element type determination
                            element_type = self._determine_pptx_element_type(shape, slide_title)
                            
                            # Extract rich text formatting if available
                            font_info = self._extract_shape_formatting(shape)
                            
                            element = DocumentElement(
                                text=shape.text.strip(),
                                page_number=slide_num,
                                element_type=element_type,
                                bbox=(shape.left, shape.top, shape.left + shape.width, shape.top + shape.height),
                                confidence=1.0,
                                language=self._detect_language(shape.text),
                                font_info=font_info
                            )
                            elements.append(element)
                        
                        # Handle tables
                        elif hasattr(shape, 'table'):
                            table_elements = self._extract_pptx_table(shape.table, slide_num)
                            elements.extend(table_elements)
                        
                        # Handle charts (basic info extraction)
                        elif hasattr(shape, 'chart'):
                            chart_info = self._extract_pptx_chart_info(shape.chart, slide_num)
                            if chart_info:
                                elements.append(chart_info)
                        
                        # Handle images with OCR if needed
                        elif shape.shape_type == 13:  # Picture shape type
                            image_element = await self._extract_pptx_image_text(shape, slide_num, file_path)
                            if image_element:
                                elements.append(image_element)
                    
                    except Exception as shape_error:
                        self.logger.warning(f"Error processing shape {shape_idx} in slide {slide_num}: {shape_error}")
                        continue
                
                # Extract slide notes if available
                if hasattr(slide, 'notes_slide') and slide.notes_slide:
                    notes_text = self._extract_slide_notes(slide.notes_slide)
                    if notes_text:
                        elements.append(DocumentElement(
                            text=notes_text,
                            page_number=slide_num,
                            element_type='notes',
                            bbox=(0, 1000, 500, 1100),
                            confidence=1.0,
                            language=self._detect_language(notes_text),
                            font_info={'type': 'slide_notes', 'slide_number': slide_num}
                        ))
            
            # Extract presentation metadata
            metadata = {
                'file_type': 'pptx', 
                'size': file_path.stat().st_size,
                'slides': len(prs.slides),
                'title': getattr(prs.core_properties, 'title', ''),
                'author': getattr(prs.core_properties, 'author', ''),
                'subject': getattr(prs.core_properties, 'subject', ''),
                'created': str(getattr(prs.core_properties, 'created', '')),
                'modified': str(getattr(prs.core_properties, 'modified', ''))
            }
            
            self.logger.info(f"PowerPoint processing completed. Extracted {len(elements)} elements from {len(prs.slides)} slides")
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=len(prs.slides),
                elements=elements,
                metadata=metadata
            )
            
        except Exception as e:
            self.logger.error(f"Error processing PowerPoint file: {e}")
            raise
    
    async def _process_image(self, file_path: Path) -> ProcessedDocument:
        """Process image files with enhanced OCR and preprocessing"""
        try:
            self.logger.info(f"Starting enhanced image processing for: {file_path.name}")
            
            # Load and validate image
            image = cv2.imread(str(file_path))
            original_image = None
            
            if image is None:
                # Try with PIL as fallback
                try:
                    pil_image = Image.open(str(file_path))
                    # Convert PIL to cv2 format
                    image = cv2.cvtColor(np.array(pil_image), cv2.COLOR_RGB2BGR)
                    original_image = np.array(pil_image)
                    self.logger.info("Image loaded using PIL fallback")
                except Exception as pil_error:
                    raise ValueError(f"Could not load image {file_path.name}. CV2 error: Image is None. PIL error: {pil_error}")
            else:
                original_image = cv2.cvtColor(image, cv2.COLOR_BGR2RGB)
            
            # Enhanced image preprocessing for better OCR
            processed_images = self._preprocess_image_for_ocr(image)
            
            elements = []
            best_results = []
            
            # Try OCR on multiple preprocessed versions and combine results
            for i, processed_img in enumerate(processed_images):
                try:
                    self.logger.info(f"Running OCR on preprocessed image {i+1}/{len(processed_images)}...")
                    
                    ocr_results = []
                    
                    # Choose OCR method based on what's available
                    if self.use_tesseract_fallback or not self.primary_ocr_reader:
                        # Use Tesseract OCR
                        ocr_results = self._extract_text_with_tesseract_direct(processed_img, i)
                    else:
                        # Use EasyOCR
                        # Save processed image temporarily for EasyOCR
                        temp_path = file_path.parent / f"temp_processed_{i}_{file_path.name}"
                        cv2.imwrite(str(temp_path), processed_img)
                        
                        try:
                            selected_reader = self.primary_ocr_reader or self.ocr_reader
                            if selected_reader:
                                # Perform OCR with EasyOCR
                                ocr_results = selected_reader.readtext(str(temp_path))
                            else:
                                self.logger.warning("No EasyOCR reader available")
                                ocr_results = []
                        except Exception as ocr_error:
                            self.logger.warning(f"EasyOCR failed: {ocr_error}, trying Tesseract...")
                            ocr_results = self._extract_text_with_tesseract_direct(processed_img, i)
                        finally:
                            # Clean up temp file
                            if temp_path.exists():
                                temp_path.unlink()
                    
                    # Store results with preprocessing info
                    for result in ocr_results:
                        if len(result) >= 3 and result[2] > 0.2:  # Lowered confidence threshold
                            best_results.append((result, i))
                    
                    self.logger.info(f"OCR on version {i+1} found {len(ocr_results)} text regions")
                    
                except Exception as ocr_error:
                    self.logger.warning(f"OCR failed on preprocessed version {i+1}: {ocr_error}")
                    continue
            
            # Also try Tesseract OCR as backup
            try:
                self.logger.info("Trying Tesseract OCR as backup...")
                tesseract_results = self._extract_text_with_tesseract(image, file_path)
                best_results.extend(tesseract_results)
            except Exception as tess_error:
                self.logger.warning(f"Tesseract OCR failed: {tess_error}")
            
            # Deduplicate and sort results by confidence
            unique_results = self._deduplicate_ocr_results(best_results)
            self.logger.info(f"Total unique OCR results: {len(unique_results)}")
            
            # Process all unique OCR results
            try:
                for i, (result_data, preprocessing_version) in enumerate(unique_results):
                    try:
                        if isinstance(result_data, tuple) and len(result_data) >= 3:
                            # EasyOCR format: (bbox_coords, text, confidence)
                            bbox_coords = result_data[0]
                            text = result_data[1]
                            confidence = result_data[2]
                        else:
                            # Tesseract format: (text, bbox, confidence)
                            text = result_data[0]
                            bbox_coords = result_data[1]
                            confidence = result_data[2]
                        
                        self.logger.debug(f"Processing OCR result {i+1}: '{text}' (confidence: {confidence:.2f}, version: {preprocessing_version})")
                        
                        if confidence > 0.25 and text.strip() and len(text.strip()) > 1:  # More lenient threshold
                            # Convert bbox format safely
                            try:
                                if isinstance(bbox_coords, list) and len(bbox_coords) > 0:
                                    # EasyOCR format
                                    if isinstance(bbox_coords[0], list):
                                        x_coords = [float(coord[0]) for coord in bbox_coords]
                                        y_coords = [float(coord[1]) for coord in bbox_coords]
                                        bbox = (min(x_coords), min(y_coords), max(x_coords), max(y_coords))
                                    else:
                                        # Tesseract format: [x, y, w, h]
                                        x, y, w, h = bbox_coords
                                        bbox = (float(x), float(y), float(x + w), float(y + h))
                                else:
                                    # Fallback bbox
                                    bbox = (0, i*30, 200, (i+1)*30)
                            except (ValueError, IndexError, TypeError) as bbox_error:
                                self.logger.warning(f"Invalid bbox for text '{text}': {bbox_error}")
                                bbox = (0, i*30, 200, (i+1)*30)
                            
                            # Enhanced element type determination
                            element_type = self._determine_text_element_type(text, confidence)
                            
                            # Detect language with improved accuracy
                            detected_lang = self._detect_language(text)
                            
                            element = DocumentElement(
                                text=text.strip(),
                                page_number=1,
                                element_type=element_type,
                                bbox=bbox,
                                confidence=float(confidence),
                                language=detected_lang,
                                font_info={'preprocessing_version': preprocessing_version, 'ocr_engine': 'easyocr' if isinstance(result_data, tuple) else 'tesseract'}
                            )
                            elements.append(element)
                            
                    except Exception as element_error:
                        self.logger.error(f"Error processing OCR result {i+1}: {element_error}")
                        continue
            
            except Exception as ocr_error:
                self.logger.error(f"OCR processing failed: {ocr_error}")
                # Create a fallback element indicating OCR failure
                elements.append(DocumentElement(
                    text=f"OCR processing failed for image {file_path.name}. Error: {str(ocr_error)}",
                    page_number=1,
                    element_type='error',
                    bbox=(0, 0, 100, 30),
                    confidence=0.0,
                    language='en'
                ))
            
            # Perform layout analysis if model is available
            if self.layout_model and len(elements) > 0:
                try:
                    self.logger.info("Running layout analysis...")
                    layout = self.layout_model.detect(image)
                    # This would enhance the element classification
                    # Implementation depends on specific layout model output
                    self.logger.info("Layout analysis completed")
                except Exception as e:
                    self.logger.warning(f"Layout analysis failed: {e}")
            
            # If no elements were extracted, add a placeholder
            if not elements:
                self.logger.warning(f"No text found in image {file_path.name}")
                elements.append(DocumentElement(
                    text=f"No readable text found in image {file_path.name}",
                    page_number=1,
                    element_type='placeholder',
                    bbox=(0, 0, 100, 30),
                    confidence=0.0,
                    language='en'
                ))
            
            self.logger.info(f"Image processing completed. Extracted {len(elements)} elements")
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=elements,
                metadata={
                    'file_type': 'image', 
                    'size': file_path.stat().st_size,
                    'ocr_elements': len([e for e in elements if e.element_type in ['text', 'handwriting', 'heading', 'number']]),
                    'image_dimensions': f"{image.shape[1]}x{image.shape[0]}" if image is not None else "unknown",
                    'preprocessing_versions': len(processed_images),
                    'total_ocr_results': len(unique_results),
                    'confidence_avg': sum(e.confidence for e in elements) / len(elements) if elements else 0
                }
            )
            
        except Exception as e:
            self.logger.error(f"Critical error in image processing for {file_path.name}: {e}")
            # Return a document with error information instead of raising
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=[DocumentElement(
                    text=f"Failed to process image {file_path.name}: {str(e)}",
                    page_number=1,
                    element_type='error',
                    bbox=(0, 0, 100, 30),
                    confidence=0.0,
                    language='en'
                )],
                metadata={'file_type': 'image', 'error': str(e), 'size': file_path.stat().st_size if file_path.exists() else 0}
            )
    
    async def _process_text(self, file_path: Path) -> ProcessedDocument:
        """Process plain text and RTF files"""
        try:
            # Detect file encoding
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                encoding = chardet.detect(raw_data)['encoding'] or 'utf-8'
            
            # Read text content
            with open(file_path, 'r', encoding=encoding) as f:
                content = f.read()
            
            elements = []
            
            # Split into paragraphs
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            
            if not paragraphs:
                # If no double newlines, split by single newlines
                paragraphs = [p.strip() for p in content.split('\n') if p.strip()]
            
            for i, paragraph in enumerate(paragraphs):
                element = DocumentElement(
                    text=paragraph,
                    page_number=1,
                    element_type='paragraph',
                    bbox=(0, i*30, 500, (i+1)*30),
                    confidence=1.0,
                    language=self._detect_language(paragraph)
                )
                elements.append(element)
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=elements,
                metadata={'file_type': 'text', 'encoding': encoding, 'size': file_path.stat().st_size}
            )
            
        except Exception as e:
            self.logger.error(f"Error processing text file: {e}")
            raise
    
    async def _process_markdown(self, file_path: Path) -> ProcessedDocument:
        """Process Markdown files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Convert markdown to HTML for better structure analysis
            html_content = markdown.markdown(content)
            soup = BeautifulSoup(html_content, 'html.parser')
            
            elements = []
            
            # Extract structured content
            for i, element in enumerate(soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'li', 'blockquote'])):
                text = element.get_text().strip()
                if text:
                    element_type = 'heading' if element.name.startswith('h') else 'paragraph'
                    
                    doc_element = DocumentElement(
                        text=text,
                        page_number=1,
                        element_type=element_type,
                        bbox=(0, i*30, 500, (i+1)*30),
                        confidence=1.0,
                        language=self._detect_language(text),
                        font_info={'tag': element.name}
                    )
                    elements.append(doc_element)
            
            # If no structured elements found, fall back to plain text processing
            if not elements:
                paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
                for i, paragraph in enumerate(paragraphs):
                    element = DocumentElement(
                        text=paragraph,
                        page_number=1,
                        element_type='paragraph',
                        bbox=(0, i*30, 500, (i+1)*30),
                        confidence=1.0,
                        language=self._detect_language(paragraph)
                    )
                    elements.append(element)
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=elements,
                metadata={'file_type': 'markdown', 'size': file_path.stat().st_size}
            )
            
        except Exception as e:
            self.logger.error(f"Error processing markdown file: {e}")
            raise
    
    async def _process_spreadsheet(self, file_path: Path) -> ProcessedDocument:
        """Process CSV and Excel files"""
        try:
            file_ext = file_path.suffix.lower()
            
            # Read the file based on format
            if file_ext == '.csv':
                df = pd.read_csv(file_path, encoding='utf-8')
            elif file_ext in ['.xlsx', '.xls']:
                if not OPENPYXL_AVAILABLE and file_ext == '.xlsx':
                    raise ValueError("openpyxl is required for Excel files")
                df = pd.read_excel(file_path)
            else:
                raise ValueError(f"Unsupported spreadsheet format: {file_ext}")
            
            elements = []
            
            # Add headers as a heading element
            headers = list(df.columns)
            header_text = ' | '.join(str(h) for h in headers)
            
            header_element = DocumentElement(
                text=header_text,
                page_number=1,
                element_type='heading',
                bbox=(0, 0, 500, 30),
                confidence=1.0,
                language=self._detect_language(header_text),
                font_info={'type': 'header_row'}
            )
            elements.append(header_element)
            
            # Process each row as a table element
            for i, (index, row) in enumerate(df.iterrows(), 1):
                row_text = ' | '.join(str(value) for value in row.values)
                if row_text.strip():
                    element = DocumentElement(
                        text=row_text,
                        page_number=1,
                        element_type='table',
                        bbox=(0, i*30, 500, (i+1)*30),
                        confidence=1.0,
                        language=self._detect_language(row_text),
                        font_info={'row_index': index}
                    )
                    elements.append(element)
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=elements,
                metadata={
                    'file_type': 'spreadsheet', 
                    'format': file_ext,
                    'rows': len(df),
                    'columns': len(df.columns),
                    'size': file_path.stat().st_size
                }
            )
            
        except Exception as e:
            self.logger.error(f"Error processing spreadsheet file: {e}")
            raise
    
    async def _process_json(self, file_path: Path) -> ProcessedDocument:
        """Process JSON files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            elements = []
            
            def flatten_json(obj, prefix='', level=0):
                """Recursively flatten JSON structure"""
                items = []
                if isinstance(obj, dict):
                    for key, value in obj.items():
                        new_key = f"{prefix}.{key}" if prefix else key
                        if isinstance(value, (dict, list)):
                            items.extend(flatten_json(value, new_key, level+1))
                        else:
                            items.append((new_key, str(value), level))
                elif isinstance(obj, list):
                    for i, item in enumerate(obj):
                        new_key = f"{prefix}[{i}]"
                        if isinstance(item, (dict, list)):
                            items.extend(flatten_json(item, new_key, level+1))
                        else:
                            items.append((new_key, str(item), level))
                else:
                    items.append((prefix, str(obj), level))
                return items
            
            flattened = flatten_json(data)
            
            for i, (key, value, level) in enumerate(flattened):
                text = f"{key}: {value}"
                element_type = 'heading' if level == 0 else 'paragraph'
                
                element = DocumentElement(
                    text=text,
                    page_number=1,
                    element_type=element_type,
                    bbox=(level*20, i*25, 500, (i+1)*25),
                    confidence=1.0,
                    language=self._detect_language(value),
                    font_info={'level': level, 'key': key}
                )
                elements.append(element)
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=elements,
                metadata={'file_type': 'json', 'size': file_path.stat().st_size}
            )
            
        except Exception as e:
            self.logger.error(f"Error processing JSON file: {e}")
            raise
    
    async def _process_xml(self, file_path: Path) -> ProcessedDocument:
        """Process XML files"""
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()
            
            elements = []
            
            def process_element(elem, level=0, index=0):
                """Recursively process XML elements"""
                text_content = elem.text.strip() if elem.text else ''
                
                if text_content:
                    element_text = f"{elem.tag}: {text_content}"
                    element_type = 'heading' if level == 0 else 'paragraph'
                    
                    element = DocumentElement(
                        text=element_text,
                        page_number=1,
                        element_type=element_type,
                        bbox=(level*20, index*25, 500, (index+1)*25),
                        confidence=1.0,
                        language=self._detect_language(text_content),
                        font_info={'tag': elem.tag, 'level': level}
                    )
                    elements.append(element)
                
                # Process attributes
                if elem.attrib:
                    for attr_name, attr_value in elem.attrib.items():
                        attr_text = f"{elem.tag}@{attr_name}: {attr_value}"
                        element = DocumentElement(
                            text=attr_text,
                            page_number=1,
                            element_type='paragraph',
                            bbox=(level*20, len(elements)*25, 500, (len(elements)+1)*25),
                            confidence=1.0,
                            language=self._detect_language(attr_value),
                            font_info={'tag': elem.tag, 'level': level, 'type': 'attribute'}
                        )
                        elements.append(element)
                
                # Process child elements
                for child in elem:
                    process_element(child, level+1, len(elements))
            
            process_element(root)
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=elements,
                metadata={'file_type': 'xml', 'root_tag': root.tag, 'size': file_path.stat().st_size}
            )
            
        except Exception as e:
            self.logger.error(f"Error processing XML file: {e}")
            raise
    
    async def _process_html(self, file_path: Path) -> ProcessedDocument:
        """Process HTML files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            soup = BeautifulSoup(content, 'html.parser')
            elements = []
            
            # Extract meaningful content elements
            content_tags = ['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'div', 'span', 'li', 'td', 'th']
            
            for i, element in enumerate(soup.find_all(content_tags)):
                text = element.get_text().strip()
                if text and len(text) > 3:  # Filter out very short text
                    element_type = 'heading' if element.name.startswith('h') else 'paragraph'
                    if element.name in ['th', 'td']:
                        element_type = 'table'
                    
                    doc_element = DocumentElement(
                        text=text,
                        page_number=1,
                        element_type=element_type,
                        bbox=(0, i*30, 500, (i+1)*30),
                        confidence=1.0,
                        language=self._detect_language(text),
                        font_info={'tag': element.name, 'class': element.get('class', [])}
                    )
                    elements.append(doc_element)
            
            # Extract title if available
            title_tag = soup.find('title')
            if title_tag and title_tag.get_text().strip():
                title_element = DocumentElement(
                    text=title_tag.get_text().strip(),
                    page_number=1,
                    element_type='heading',
                    bbox=(0, -30, 500, 0),  # Position before other content
                    confidence=1.0,
                    language=self._detect_language(title_tag.get_text()),
                    font_info={'tag': 'title'}
                )
                elements.insert(0, title_element)
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=elements,
                metadata={'file_type': 'html', 'size': file_path.stat().st_size}
            )
            
        except Exception as e:
            self.logger.error(f"Error processing HTML file: {e}")
            raise
    
    async def _process_odt(self, file_path: Path) -> ProcessedDocument:
        """Process OpenDocument Text files"""
        try:
            if not ODT_AVAILABLE:
                raise ValueError("odfpy is required for ODT files")
            
            doc = load(file_path)
            elements = []
            
            # Extract all paragraph elements
            paragraphs = doc.getElementsByType(P)
            
            for i, paragraph in enumerate(paragraphs):
                text = teletype.extractText(paragraph).strip()
                if text:
                    element = DocumentElement(
                        text=text,
                        page_number=1,
                        element_type='paragraph',
                        bbox=(0, i*30, 500, (i+1)*30),
                        confidence=1.0,
                        language=self._detect_language(text)
                    )
                    elements.append(element)
            
            return ProcessedDocument(
                filename=file_path.name,
                total_pages=1,
                elements=elements,
                metadata={'file_type': 'odt', 'size': file_path.stat().st_size}
            )
            
        except Exception as e:
            self.logger.error(f"Error processing ODT file: {e}")
            raise
    
    def _detect_language(self, text: str) -> str:
        """Enhanced language detection based on character analysis with Indian language support"""
        if not text or not text.strip():
            return 'unknown'
        
        # Clean and normalize the text first
        try:
            # Ensure proper Unicode handling
            if isinstance(text, bytes):
                text = text.decode('utf-8', errors='ignore')
            
            # Remove control characters and normalize whitespace
            import re
            text = re.sub(r'[\x00-\x1f\x7f-\x9f]', '', text)
            text = re.sub(r'\s+', ' ', text).strip()
            
            if len(text) < 3:
                return 'unknown'
            
        except Exception as e:
            self.logger.debug(f"Text normalization failed: {e}")
            return 'unknown'
        
        # Try using the advanced Indian language detector first
        try:
            from src.utils.indian_language_detector import detect_indian_language
            detection = detect_indian_language(text)
            if detection.confidence > 0.4:  # Lowered threshold for better detection
                self.logger.debug(f"Detected language: {detection.language_code} ({detection.confidence:.2f})")
                return detection.language_code
        except Exception as e:
            self.logger.debug(f"Indian language detection failed, falling back to simple detection: {e}")
        
        # Fallback: Count different script characters
        latin_chars = sum(1 for c in text if c.isascii() and c.isalpha())
        arabic_chars = sum(1 for c in text if '\u0600' <= c <= '\u06FF')
        # Devanagari script (Hindi, Marathi, etc.)
        devanagari_chars = sum(1 for c in text if '\u0900' <= c <= '\u097F')
        # Kannada script
        kannada_chars = sum(1 for c in text if '\u0C80' <= c <= '\u0CFF')
        # Telugu script
        telugu_chars = sum(1 for c in text if '\u0C00' <= c <= '\u0C7F')
        # Tamil script
        tamil_chars = sum(1 for c in text if '\u0B80' <= c <= '\u0BFF')
        # Bengali script
        bengali_chars = sum(1 for c in text if '\u0980' <= c <= '\u09FF')
        # Gujarati script
        gujarati_chars = sum(1 for c in text if '\u0A80' <= c <= '\u0AFF')
        # Malayalam script
        malayalam_chars = sum(1 for c in text if '\u0D00' <= c <= '\u0D7F')
        # Punjabi script (Gurmukhi)
        punjabi_chars = sum(1 for c in text if '\u0A00' <= c <= '\u0A7F')
        
        total_chars = (latin_chars + arabic_chars + devanagari_chars + 
                      kannada_chars + telugu_chars + tamil_chars + 
                      bengali_chars + gujarati_chars + malayalam_chars + punjabi_chars)
        
        if total_chars == 0:
            return 'unknown'
        
        # Determine dominant script with threshold of 20%
        threshold = 0.2
        
        if kannada_chars / total_chars > threshold:
            return 'kn'
        elif telugu_chars / total_chars > threshold:
            return 'te'
        elif tamil_chars / total_chars > threshold:
            return 'ta'
        elif devanagari_chars / total_chars > threshold:
            return 'hi'  # Could be Hindi or Marathi, default to Hindi
        elif bengali_chars / total_chars > threshold:
            return 'bn'
        elif gujarati_chars / total_chars > threshold:
            return 'gu'
        elif malayalam_chars / total_chars > threshold:
            return 'ml'
        elif punjabi_chars / total_chars > threshold:
            return 'pa'
        elif arabic_chars / total_chars > threshold:
            return 'ar'
        else:
            return 'en'
    
    def _preprocess_image_for_ocr(self, image: np.ndarray) -> List[np.ndarray]:
        """Apply various preprocessing techniques to improve OCR accuracy"""
        preprocessed_images = []
        
        try:
            # Original image
            preprocessed_images.append(image.copy())
            
            # Convert to grayscale
            gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
            
            # 1. Enhance contrast using CLAHE
            clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
            enhanced = clahe.apply(gray)
            preprocessed_images.append(cv2.cvtColor(enhanced, cv2.COLOR_GRAY2BGR))
            
            # 2. Gaussian blur + threshold
            blurred = cv2.GaussianBlur(gray, (5, 5), 0)
            _, thresh1 = cv2.threshold(blurred, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            preprocessed_images.append(cv2.cvtColor(thresh1, cv2.COLOR_GRAY2BGR))
            
            # 3. Adaptive threshold
            adaptive_thresh = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2)
            preprocessed_images.append(cv2.cvtColor(adaptive_thresh, cv2.COLOR_GRAY2BGR))
            
            # 4. Morphological operations to clean up
            kernel = np.ones((2,2), np.uint8)
            morphed = cv2.morphologyEx(thresh1, cv2.MORPH_CLOSE, kernel)
            preprocessed_images.append(cv2.cvtColor(morphed, cv2.COLOR_GRAY2BGR))
            
            # 5. Noise removal
            denoised = cv2.fastNlMeansDenoising(gray)
            preprocessed_images.append(cv2.cvtColor(denoised, cv2.COLOR_GRAY2BGR))
            
            self.logger.info(f"Generated {len(preprocessed_images)} preprocessed versions")
            
        except Exception as e:
            self.logger.warning(f"Error in image preprocessing: {e}")
            # Return at least the original image
            if not preprocessed_images:
                preprocessed_images = [image.copy()]
        
        return preprocessed_images
    
    def _extract_text_with_tesseract(self, image: np.ndarray, file_path: Path) -> List[Tuple[Any, int]]:
        """Extract text using Tesseract OCR as backup"""
        results = []
        
        try:
            # Convert image to PIL format for Tesseract
            pil_image = Image.fromarray(cv2.cvtColor(image, cv2.COLOR_BGR2RGB))
            
            # Get text with bounding boxes
            data = pytesseract.image_to_data(pil_image, output_type=pytesseract.Output.DICT)
            
            for i in range(len(data['text'])):
                text = data['text'][i].strip()
                confidence = int(data['conf'][i])
                
                if text and confidence > 30:  # Minimum confidence
                    x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
                    bbox = [x, y, w, h]
                    
                    # Format: (text, bbox, confidence/100)
                    results.append(((text, bbox, confidence/100.0), -1))  # -1 indicates Tesseract
            
            self.logger.info(f"Tesseract found {len(results)} text regions")
            
        except Exception as e:
            self.logger.warning(f"Tesseract processing failed: {e}")
        
        return results
    
    def _extract_text_with_tesseract_direct(self, image: np.ndarray, version: int) -> List[Tuple[Any, int]]:
        """Extract text using Tesseract OCR directly from image array with multilingual support"""
        results = []
        
        try:
            # Convert image to PIL format for Tesseract
            pil_image = Image.fromarray(cv2.cvtColor(image, cv2.COLOR_BGR2RGB))
            
            # Configure Tesseract for multiple languages
            # Try different language combinations based on common usage
            lang_configs = [
                'eng+hin+kan',  # English + Hindi + Kannada
                'eng+hin',      # English + Hindi
                'eng+kan',      # English + Kannada
                'eng',          # English only
            ]
            
            best_results = []
            best_count = 0
            
            for lang_config in lang_configs:
                try:
                    # Get text with bounding boxes using current language config
                    data = pytesseract.image_to_data(
                        pil_image, 
                        output_type=pytesseract.Output.DICT,
                        lang=lang_config,
                        config='--psm 6'  # Assume a single uniform block of text
                    )
                    
                    current_results = []
                    for i in range(len(data['text'])):
                        text = data['text'][i].strip()
                        confidence = int(data['conf'][i])
                        
                        if text and confidence > 25:  # Lower threshold for multilingual
                            x, y, w, h = data['left'][i], data['top'][i], data['width'][i], data['height'][i]
                            bbox = [x, y, w, h]
                            
                            # Format: (text, bbox, confidence/100)
                            current_results.append((text, bbox, confidence/100.0))
                    
                    # Keep the configuration that found the most text
                    if len(current_results) > best_count:
                        best_count = len(current_results)
                        best_results = current_results
                        self.logger.debug(f"Best results so far with {lang_config}: {len(current_results)} regions")
                    
                    # If we found good results, don't try more configs
                    if len(current_results) > 5:
                        break
                        
                except Exception as lang_error:
                    self.logger.debug(f"Language config {lang_config} failed: {lang_error}")
                    continue
            
            # Format results for compatibility with EasyOCR format
            for text, bbox, confidence in best_results:
                results.append(((text, bbox, confidence), version))
            
            self.logger.info(f"Tesseract direct processing found {len(results)} text regions")
            
        except Exception as e:
            self.logger.warning(f"Tesseract direct processing failed: {e}")
        
        return results
    
    def _deduplicate_ocr_results(self, results: List[Tuple[Any, int]]) -> List[Tuple[Any, int]]:
        """Remove duplicate OCR results based on text similarity and position"""
        unique_results = []
        
        for current_result, version in results:
            is_duplicate = False
            
            # Extract text for comparison
            if isinstance(current_result, tuple) and len(current_result) >= 3:
                current_text = current_result[1]  # EasyOCR format
            else:
                current_text = current_result[0]  # Tesseract format
            
            # Check against existing results
            for existing_result, _ in unique_results:
                if isinstance(existing_result, tuple) and len(existing_result) >= 3:
                    existing_text = existing_result[1]
                else:
                    existing_text = existing_result[0]
                
                # Simple similarity check
                if self._text_similarity(current_text, existing_text) > 0.8:
                    is_duplicate = True
                    break
            
            if not is_duplicate:
                unique_results.append((current_result, version))
        
        # Sort by confidence (highest first)
        try:
            unique_results.sort(key=lambda x: self._get_confidence(x[0]), reverse=True)
        except:
            pass
        
        return unique_results
    
    def _text_similarity(self, text1: str, text2: str) -> float:
        """Calculate simple text similarity ratio"""
        if not text1 or not text2:
            return 0.0
        
        # Simple character-based similarity
        text1_clean = text1.lower().strip()
        text2_clean = text2.lower().strip()
        
        if text1_clean == text2_clean:
            return 1.0
        
        # Calculate overlap ratio
        shorter = min(len(text1_clean), len(text2_clean))
        longer = max(len(text1_clean), len(text2_clean))
        
        if shorter == 0:
            return 0.0
        
        # Count common characters in order
        common = 0
        for i in range(shorter):
            if i < len(text1_clean) and i < len(text2_clean) and text1_clean[i] == text2_clean[i]:
                common += 1
        
        return common / longer
    
    def _get_confidence(self, result: Any) -> float:
        """Extract confidence from OCR result"""
        try:
            if isinstance(result, tuple) and len(result) >= 3:
                return float(result[2])  # EasyOCR or Tesseract format
            return 0.0
        except:
            return 0.0
    
    def _determine_text_element_type(self, text: str, confidence: float) -> str:
        """Determine the type of text element based on content and confidence"""
        text_clean = text.strip().lower()
        
        # Check for numbers
        if text_clean.replace('.', '').replace(',', '').replace('-', '').isdigit():
            return 'number'
        
        # Check for headings (short, often capitalized)
        if len(text_clean) < 50 and (text.isupper() or text.istitle()):
            return 'heading'
        
        # Check for handwriting (low confidence)
        if confidence < 0.6:
            return 'handwriting'
        
        # Default to text
        return 'text'
    
    def get_document_stats(self, document: ProcessedDocument) -> Dict[str, Any]:
        """Generate statistics about the processed document"""
        stats = {
            'total_elements': len(document.elements),
            'total_pages': document.total_pages,
            'element_types': {},
            'languages': {},
            'avg_confidence': 0,
            'total_text_length': 0
        }
        
        for element in document.elements:
            # Count element types
            stats['element_types'][element.element_type] = \
                stats['element_types'].get(element.element_type, 0) + 1
            
            # Count languages
            if element.language:
                stats['languages'][element.language] = \
                    stats['languages'].get(element.language, 0) + 1
            
            # Calculate averages
            stats['avg_confidence'] += element.confidence
            stats['total_text_length'] += len(element.text)
        
        if document.elements:
            stats['avg_confidence'] /= len(document.elements)
        
        return stats
    
    # Helper methods for PowerPoint processing
    def _extract_slide_title(self, slide):
        """Extract slide title if available"""
        try:
            if hasattr(slide, 'shapes'):
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        # First text shape is often the title
                        return shape.text.strip()
            return None
        except Exception:
            return None
    
    def _determine_pptx_element_type(self, shape, slide_title):
        """Determine element type for PowerPoint shapes"""
        try:
            if hasattr(shape, 'text') and shape.text:
                text = shape.text.strip()
                # If this is the first text and matches slide title, it's a heading
                if slide_title and text == slide_title:
                    return 'heading'
                # Short text in uppercase might be a heading
                elif len(text) < 50 and text.isupper():
                    return 'heading'
                else:
                    return 'paragraph'
            return 'text'
        except Exception:
            return 'text'
    
    def _extract_shape_formatting(self, shape):
        """Extract formatting information from shape"""
        try:
            formatting = {'shape_type': 'text'}
            if hasattr(shape, 'text_frame'):
                formatting['has_text_frame'] = True
            return formatting
        except Exception:
            return {'shape_type': 'unknown'}
    
    def _extract_pptx_table(self, table, slide_num):
        """Extract table content from PowerPoint"""
        elements = []
        try:
            for row_idx, row in enumerate(table.rows):
                row_text = []
                for cell in row.cells:
                    row_text.append(cell.text.strip() if cell.text else '')
                
                if any(text for text in row_text):  # If any cell has content
                    element = DocumentElement(
                        text=' | '.join(row_text),
                        page_number=slide_num,
                        element_type='table',
                        bbox=(0, row_idx * 30, 500, (row_idx + 1) * 30),
                        confidence=1.0,
                        language=self._detect_language(' '.join(row_text))
                    )
                    elements.append(element)
        except Exception as e:
            self.logger.warning(f"Error extracting table: {e}")
        return elements
    
    def _extract_pptx_chart_info(self, chart, slide_num):
        """Extract basic chart information"""
        try:
            chart_text = f"Chart: {getattr(chart, 'chart_title', 'Untitled Chart')}"
            return DocumentElement(
                text=chart_text,
                page_number=slide_num,
                element_type='figure',
                bbox=(0, 0, 200, 50),
                confidence=1.0,
                language='en'
            )
        except Exception:
            return None
    
    async def _extract_pptx_image_text(self, shape, slide_num, file_path):
        """Extract text from images in PowerPoint using OCR"""
        try:
            # For now, just return a placeholder - would need image extraction
            return DocumentElement(
                text="Image content (OCR not implemented for embedded images)",
                page_number=slide_num,
                element_type='image',
                bbox=(0, 0, 100, 50),
                confidence=0.5,
                language='en'
            )
        except Exception:
            return None
    
    def _extract_slide_notes(self, notes_slide):
        """Extract notes from slide"""
        try:
            if hasattr(notes_slide, 'notes_text_frame') and notes_slide.notes_text_frame:
                return notes_slide.notes_text_frame.text.strip()
            return None
        except Exception:
            return None
    
    def _extract_text_from_binary_ppt(self, file_path: Path) -> List[str]:
        """Extract readable text strings from binary PPT files"""
        text_content = []
        
        try:
            with open(file_path, 'rb') as f:
                # Read file in chunks
                chunk_size = 8192
                content = b''
                
                while True:
                    chunk = f.read(chunk_size)
                    if not chunk:
                        break
                    content += chunk
                    
                    # Limit total content to avoid memory issues
                    if len(content) > 10 * 1024 * 1024:  # 10MB limit
                        break
            
            # Look for printable ASCII strings
            import re
            
            # Find strings of printable ASCII characters (minimum 4 characters)
            ascii_strings = re.findall(rb'[\x20-\x7E]{4,}', content)
            
            # Convert to text and filter
            for string_bytes in ascii_strings:
                try:
                    text = string_bytes.decode('ascii', errors='ignore').strip()
                    
                    # Filter out obvious binary artifacts and common PPT metadata
                    if (len(text) > 3 and 
                        not text.startswith(('\\x', 'Microsoft', 'PowerPoint', 'Times New Roman')) and
                        not all(c in '0123456789.-+' for c in text) and  # Skip pure numbers
                        not text.isupper() or len(text) < 20):  # Skip long uppercase strings (likely metadata)
                        
                        # Basic content filtering
                        if any(keyword in text.lower() for keyword in 
                               ['slide', 'title', 'content', 'text', 'bullet', 'point']):
                            text_content.append(text)
                        elif len([c for c in text if c.isalpha()]) > len(text) * 0.7:  # Mostly letters
                            text_content.append(text)
                            
                except UnicodeDecodeError:
                    continue
            
            # Also try UTF-16 encoding (common in Office files)
            try:
                utf16_strings = re.findall(rb'[\x20-\x7E\x00]{8,}', content)
                for string_bytes in utf16_strings[:50]:  # Limit to first 50 matches
                    try:
                        text = string_bytes.decode('utf-16le', errors='ignore').strip()
                        if len(text) > 3 and text.isprintable():
                            text_content.append(text)
                    except UnicodeDecodeError:
                        continue
            except Exception:
                pass
            
            # Remove duplicates and sort by length (longer strings first)
            unique_texts = list(set(text_content))
            unique_texts.sort(key=len, reverse=True)
            
            # Return top 20 text strings
            return unique_texts[:20]
            
        except Exception as e:
            self.logger.warning(f"Binary text extraction failed: {e}")
            return []
